"""
Presentation generator for financial analysis slide decks.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Optional
from datetime import datetime
import os


class PresentationGenerator:
    """
    Generate PowerPoint presentations for financial analysis.
    """
    
    def __init__(self):
        """Initialize presentation generator."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str, subtitle: Optional[str] = None):
        """
        Add title slide.
        
        Args:
            title: Main title
            subtitle: Optional subtitle
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
    
    def add_content_slide(self, title: str, content: List[str]):
        """
        Add content slide with bullet points.
        
        Args:
            title: Slide title
            content: List of bullet points
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        tf = body_shape.text_frame
        tf.text = content[0]
        
        for bullet in content[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    
    def add_chart_slide(self, title: str, chart_path: Optional[str] = None):
        """
        Add slide with chart placeholder.
        
        Args:
            title: Slide title
            chart_path: Optional path to chart image
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if chart_path and os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(1), Inches(2), 
                                   width=Inches(8), height=Inches(5))
    
    def add_table_slide(self, title: str, data: Dict):
        """
        Add slide with data table.
        
        Args:
            title: Slide title
            data: Dictionary with table data
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), 
                                              Inches(9), Inches(0.5))
        title_frame = title_shape.text_frame
        title_frame.text = title
        
        # Create table
        rows = len(data) + 1
        cols = 2
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        table.cell(0, 0).text = 'Metric'
        table.cell(0, 1).text = 'Value'
        
        # Data rows
        for i, (key, value) in enumerate(data.items(), 1):
            table.cell(i, 0).text = str(key)
            table.cell(i, 1).text = str(value)
    
    def save(self, filename: str):
        """
        Save presentation.
        
        Args:
            filename: Output filename
        """
        self.prs.save(filename)


def create_sample_presentation(output_path: str = 'sample_presentation.pptx'):
    """
    Create a sample financial analysis presentation.
    
    Args:
        output_path: Output file path
    """
    generator = PresentationGenerator()
    
    # Title slide
    generator.add_title_slide(
        'Financial Analysis Presentation',
        f'Generated: {datetime.now().strftime("%Y-%m-%d")}'
    )
    
    # Executive Summary
    generator.add_content_slide(
        'Executive Summary',
        [
            'Key findings and recommendations',
            'Market overview and conditions',
            'Valuation summary',
            'Risk assessment'
        ]
    )
    
    # Market Overview
    generator.add_content_slide(
        'Market Overview',
        [
            'Current market conditions',
            'Key economic indicators',
            'Sector performance',
            'Market outlook'
        ]
    )
    
    # Valuation
    generator.add_table_slide(
        'Valuation Summary',
        {
            'Base Case': '$100.00',
            'Bull Case': '$120.00',
            'Bear Case': '$80.00',
            'Current Price': '$95.00'
        }
    )
    
    # Risk Analysis
    generator.add_content_slide(
        'Risk Analysis',
        [
            'VaR (95%): -5.2%',
            'CVaR (95%): -7.8%',
            'Maximum Drawdown: -15.3%',
            'Beta: 1.2'
        ]
    )
    
    # Recommendations
    generator.add_content_slide(
        'Recommendations',
        [
            'Target price: $105-110',
            'Recommendation: Buy',
            'Key risks to monitor',
            'Catalysts and timeline'
        ]
    )
    
    generator.save(output_path)
    return output_path
